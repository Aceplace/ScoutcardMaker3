from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt

from scoutcardmaker.Utils import INVALID_POSITION

#Constants for wide view
CENTER_X_POS = Cm(13.0)
CENTER_Y_POS = Cm(11.5)
CENTER_Y_POS_OFF = Cm(8.7)
HORIZONTAL_COORDINATE_SIZE = Cm(0.23)
VERTICAL_COORDINATE_SIZE = Cm(0.55)
PLAYER_WIDTH = Cm(0.65)
PLAYER_HEIGHT = Cm(0.5)
HASH_SIZE = Cm(0.1)
TITLE_LEFT = 106531
TITLE_TOP = 132963
TITLE_WIDTH = 8904303
TITLE_HEIGHT = 369332
DEFENDER_WIDTH = Pt(28)
DEFENDER_HEIGHT = Pt(28)
#constants for tight view
TIGHT_CENTER_X_POS = Cm(13.0)
TIGHT_CENTER_Y_POS = Cm(11.0)
TIGHT_CENTER_Y_POS_OFF = Cm(9.8)
TIGHT_HORIZONTAL_COORDINATE_SIZE = Cm(0.46)
TIGHT_VERTICAL_COORDINATE_SIZE = Cm(1.05)
TIGHT_PLAYER_WIDTH = Cm(1.3)
TIGHT_PLAYER_HEIGHT = Cm(1.0)
TIGHT_DEFENDER_WIDTH = Pt(40)
TIGHT_DEFENDER_HEIGHT = Pt(40)
#Derived Constant Values
LEFT_SIDELINE = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 53
RIGHT_SIDELINE = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 53
LEFT_HASH = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 18
RIGHT_HASH = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 18
LEFT_TOP_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 35
LEFT_BOTTOM_OF_NUMBERS = CENTER_X_POS - HORIZONTAL_COORDINATE_SIZE * 39
RIGHT_TUP_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 35
RIGHT_BOTTOM_OF_NUMBERS = CENTER_X_POS + HORIZONTAL_COORDINATE_SIZE * 39
FIVE_YARDS = VERTICAL_COORDINATE_SIZE * 5


def player_coordinates_to_powerpoint(player_x, player_y, is_tight_view, is_for_offense):
    if is_for_offense:
        player_x *= -1
        player_y *= -1
        if is_tight_view:
            center_y_pos = TIGHT_CENTER_Y_POS_OFF
        else:
            center_y_pos = CENTER_Y_POS_OFF
    elif is_tight_view:
        center_y_pos = TIGHT_CENTER_Y_POS
    else:
        center_y_pos = CENTER_Y_POS


    if not is_tight_view:
        return (CENTER_X_POS + player_x * HORIZONTAL_COORDINATE_SIZE, center_y_pos + player_y * VERTICAL_COORDINATE_SIZE)
    else:
        return (TIGHT_CENTER_X_POS + player_x * TIGHT_HORIZONTAL_COORDINATE_SIZE, center_y_pos + player_y * TIGHT_VERTICAL_COORDINATE_SIZE)


def export_to_powerpoint(output_filename, plays, offense_library, defense_library, is_for_offense):
    presentation = Presentation()

    #do wide versions first
    for play in plays:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Personnel"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Defense"]} --- {play["Note"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_wide_formation_and_defense_slide(play, slide, offense_library, defense_library, is_for_offense)

    # do tight versions after
    for play in plays:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Personnel"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Defense"]} --- {play["Note"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_tight_formation_and_defense_slide(play, slide, offense_library, defense_library, is_for_offense)

    presentation.save(output_filename)

def export_to_powerpoint_alternating(output_filename, plays, offense_library, defender_library, is_for_offense):
    presentation = Presentation()

    for play in plays:
        # do wide versions first
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Personnel"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Defense"]} --- {play["Note"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_wide_formation_and_defense_slide(play, slide, offense_library, defender_library, is_for_offense)

        # then tight
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Number"]} {play["Personnel"]} {play["Hash"]} {play["Dnd"]} {play["Formation"]} {play["Play"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)

        text_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP + Pt(24) * 2, TITLE_WIDTH, TITLE_HEIGHT)
        text_box.text_frame.text = f'{play["Defense"]} --- {play["Note"]}'
        text_box.text_frame.paragraphs[0].font.size = Pt(24)
        add_tight_formation_and_defense_slide(play, slide, offense_library, defender_library, is_for_offense)

    presentation.save(output_filename)


def add_wide_formation_and_defense_slide(play, slide, offense_library, defense_library, is_for_offense):
    #draw sideline
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)
    slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_SIDELINE, CENTER_Y_POS - FIVE_YARDS * 3, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * 2)

    for num in range(-3, 3):
        #draw lines that go accross field at 5 yard intervals
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num, RIGHT_SIDELINE, CENTER_Y_POS + FIVE_YARDS * num)
        #draw hash marks
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_HASH, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        #draw ticks for the top of numbers and bottom of numbers
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, LEFT_TOP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_BOTTOM_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)
        slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS - HASH_SIZE, RIGHT_TUP_OF_NUMBERS, CENTER_Y_POS + num * FIVE_YARDS + HASH_SIZE)

    formation_name = play['Card Maker Formation']
    defense_name = play['Card Maker Defense']

    if formation_name:
        subformation, error_message = offense_library.get_composite_subformation(play['Hash'], formation_name)
        if not subformation:
            return
        for tag, player in subformation.players.items():
            try:
                label = offense_library.label_mappers[play['Personnel']].get_label(tag)
            except KeyError:
                label = offense_library.label_mappers['default'].get_label(tag)
            x, y = player_coordinates_to_powerpoint(player.x, player.y, False, is_for_offense)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - PLAYER_WIDTH / 2, y - PLAYER_HEIGHT / 2, PLAYER_WIDTH, PLAYER_HEIGHT)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(1.0)
            shape.text_frame.text = label if len(label) != 2 else label[1]
            shape.text_frame.paragraphs[0].font.size = Pt(12)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    if defense_name:
        composite_defense = defense_library.get_composite_defense(defense_name)
        composite_defense.place_defenders(subformation)
        for tag, defender in composite_defense.players.items():
            if (defender.placed_x, defender.placed_y) == INVALID_POSITION:
                continue
            label = defense_library.label_mappers['default'].get_label(tag)
            x, y = player_coordinates_to_powerpoint(defender.placed_x, defender.placed_y * -1, False, is_for_offense)
            text_box = slide.shapes.add_textbox(x - DEFENDER_WIDTH / 2, y - DEFENDER_HEIGHT / 2, DEFENDER_WIDTH, DEFENDER_HEIGHT)
            text_box.text_frame.text = label
            text_box.text_frame.paragraphs[0].font.size = Pt(24)


def add_tight_formation_and_defense_slide(play, slide, offense_library, defense_library, is_for_offense):
    formation_name = play['Card Maker Formation']
    defense_name = play['Card Maker Defense']

    formation = None
    x_min = -25
    x_max = 25
    y_max = 8
    offset = 0
    if formation_name:
        subformation, error_message = offense_library.get_composite_subformation(play['Hash'], formation_name)
        if not subformation:
            return
        if subformation.hash_mark == 'LT':
            x_min -= 18
            x_max -= 18
            offset = 18
        elif subformation.hash_mark == 'RT':
            x_min += 18
            x_max += 18
            offset = -18
        for tag, player in subformation.players.items():
            if player.x < x_min or player.x > x_max:
                continue
            x, y = player_coordinates_to_powerpoint(player.x + offset, player.y, True, is_for_offense)
            try:
                label = offense_library.label_mappers[play['Personnel']].get_label(tag)
            except KeyError:
                label = offense_library.label_mappers['default'].get_label(tag)
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x - TIGHT_PLAYER_WIDTH / 2, y - TIGHT_PLAYER_HEIGHT / 2,
                                           TIGHT_PLAYER_WIDTH, TIGHT_PLAYER_HEIGHT)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(1.0)
            shape.text_frame.text = label if len(label) != 2 else label[1]
            shape.text_frame.paragraphs[0].font.size = Pt(24)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    if defense_name:
        composite_defense = defense_library.get_composite_defense(defense_name)
        composite_defense.place_defenders(subformation)
        for tag, defender in composite_defense.players.items():
            if defender.placed_x < x_min or defender.placed_x > x_max or defender.placed_y > y_max:
                continue
            x, y = player_coordinates_to_powerpoint(defender.placed_x + offset, defender.placed_y * -1, True, is_for_offense)
            label = defense_library.label_mappers['default'].get_label(tag)
            text_box = slide.shapes.add_textbox(x - TIGHT_DEFENDER_WIDTH / 2, y - TIGHT_DEFENDER_HEIGHT / 2, TIGHT_DEFENDER_WIDTH,
                                                TIGHT_DEFENDER_HEIGHT)
            text_box.text_frame.text = label
            text_box.text_frame.paragraphs[0].font.size = Pt(36)

















